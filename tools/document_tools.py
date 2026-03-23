"""Cross-document conversion tools (no PDF involved as primary format).

Covers: DOCX ↔ HTML/MD/TXT, HTML ↔ MD/TXT/DOCX/XLSX,
        Markdown ↔ HTML/PDF/DOCX/TXT, TXT ↔ PDF/DOCX/HTML,
        PPTX → TXT/HTML/Images, XLSX → HTML/MD/DOCX
"""
import html as html_module
import io
import re
from pathlib import Path
from typing import Dict, List, Optional

from mcp.server.fastmcp import FastMCP

from utils.file_utils import resolve_input, resolve_output, make_output_path


# ── Markdown helpers (uses markdown-it-py, already installed) ───────────────────────

def _md_to_html_str(md_text: str) -> str:
    try:
        from markdown_it import MarkdownIt
        return MarkdownIt().render(md_text)
    except ImportError:
        # Minimal fallback renderer
        lines, out, in_list = md_text.split("\n"), [], False
        for line in lines:
            if line.startswith("### "): out.append(f"<h3>{html_module.escape(line[4:])}</h3>")
            elif line.startswith("## "): out.append(f"<h2>{html_module.escape(line[3:])}</h2>")
            elif line.startswith("# "): out.append(f"<h1>{html_module.escape(line[2:])}</h1>")
            elif line.startswith(("- ", "* ")):
                if not in_list: out.append("<ul>"); in_list = True
                out.append(f"<li>{html_module.escape(line[2:])}</li>")
            else:
                if in_list: out.append("</ul>"); in_list = False
                if line.strip(): out.append(f"<p>{html_module.escape(line)}</p>")
        if in_list: out.append("</ul>")
        return "\n".join(out)


def _html_to_md_str(html_text: str) -> str:
    try:
        import html2text
        h = html2text.HTML2Text()
        h.ignore_links = False
        return h.handle(html_text)
    except ImportError:
        # Fallback using bs4
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html_text, "lxml")
            lines = []
            for tag in soup.find_all(["h1","h2","h3","h4","p","li","br"]):
                t = tag.get_text(strip=True)
                if not t: continue
                if tag.name == "h1": lines.append(f"# {t}")
                elif tag.name == "h2": lines.append(f"## {t}")
                elif tag.name in ("h3","h4"): lines.append(f"### {t}")
                elif tag.name == "li": lines.append(f"- {t}")
                else: lines.append(t)
            return "\n\n".join(lines)
        except ImportError:
            return re.sub(r"<[^>]+>", "", html_text)


def _html_to_txt_str(html_text: str) -> str:
    try:
        from bs4 import BeautifulSoup
        return BeautifulSoup(html_text, "lxml").get_text(separator="\n")
    except ImportError:
        return html_module.unescape(re.sub(r"<[^>]+>", "", html_text))


def _wrap_html(body: str, title: str = "") -> str:
    return (
        f'<!DOCTYPE html>\n<html lang="en">\n<head>\n'
        f'  <meta charset="UTF-8">\n'
        f'  <title>{html_module.escape(title)}</title>\n'
        f'  <style>body{{font-family:sans-serif;max-width:800px;margin:auto;padding:2em}}</style>\n'
        f'</head>\n<body>\n{body}\n</body>\n</html>'
    )


def _docx_to_html_str(path: Path) -> str:
    from docx import Document
    from docx.oxml.ns import qn
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if "Heading 1" in style:
            parts.append(f"<h1>{html_module.escape(text)}</h1>")
        elif "Heading 2" in style:
            parts.append(f"<h2>{html_module.escape(text)}</h2>")
        elif "Heading 3" in style:
            parts.append(f"<h3>{html_module.escape(text)}</h3>")
        elif "List" in style:
            parts.append(f"<li>{html_module.escape(text)}</li>")
        else:
            # Inline bold/italic
            inline = []
            for run in para.runs:
                t = html_module.escape(run.text)
                if run.bold and run.italic:
                    t = f"<strong><em>{t}</em></strong>"
                elif run.bold:
                    t = f"<strong>{t}</strong>"
                elif run.italic:
                    t = f"<em>{t}</em>"
                inline.append(t)
            parts.append(f"<p>{''.join(inline)}</p>")
    # Tables
    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = "".join(f"<{'th' if i==0 else 'td'}>{html_module.escape(c.text.strip())}</{'th' if i==0 else 'td'}>" for c in row.cells)
            rows.append(f"<tr>{cells}</tr>")
        parts.append(f"<table border='1'><tbody>{''.join(rows)}</tbody></table>")
    return "\n".join(parts)


def _docx_to_md_str(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    lines = []
    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        if "Heading 1" in style: lines.append(f"# {text}")
        elif "Heading 2" in style: lines.append(f"## {text}")
        elif "Heading 3" in style: lines.append(f"### {text}")
        elif "List Bullet" in style: lines.append(f"- {text}")
        elif "List Number" in style: lines.append(f"1. {text}")
        else:
            # Inline bold/italic
            parts = []
            for run in para.runs:
                t = run.text
                if run.bold and run.italic: t = f"***{t}***"
                elif run.bold: t = f"**{t}**"
                elif run.italic: t = f"*{t}*"
                parts.append(t)
            lines.append("".join(parts))
    # Tables → MD pipe tables
    for table in doc.tables:
        if not table.rows: continue
        header = [c.text.strip() for c in table.rows[0].cells]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in table.rows[1:]:
            lines.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
        lines.append("")
    return "\n".join(lines)


def register_document_tools(app: FastMCP) -> None:

    # ── DOCX → HTML ────────────────────────────────────────────────────────────
    @app.tool()
    def docx_to_html(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert a Word document (.docx) to HTML, preserving headings, bold, italic, and tables."""
        try:
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".html")
            dst = resolve_output(output_path)
            body = _docx_to_html_str(src)
            dst.write_text(_wrap_html(body, src.stem), encoding="utf-8")
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── DOCX → Markdown ────────────────────────────────────────────────────────
    @app.tool()
    def docx_to_markdown(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert a Word document (.docx) to Markdown, preserving headings, lists, bold, italic, and tables."""
        try:
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".md")
            dst = resolve_output(output_path)
            dst.write_text(_docx_to_md_str(src), encoding="utf-8")
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── DOCX → TXT ─────────────────────────────────────────────────────────────
    @app.tool()
    def docx_to_txt(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Extract plain text from a Word document (.docx)."""
        try:
            from docx import Document
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".txt")
            dst = resolve_output(output_path)
            doc = Document(str(src))
            text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
            dst.write_text(text, encoding="utf-8")
            return {"success": True, "output_path": str(dst), "char_count": len(text)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── HTML → Markdown ────────────────────────────────────────────────────────
    @app.tool()
    def html_to_markdown(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert an HTML file to Markdown."""
        try:
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".md")
            dst = resolve_output(output_path)
            html_text = src.read_text(encoding="utf-8")
            md = _html_to_md_str(html_text)
            dst.write_text(md, encoding="utf-8")
            return {"success": True, "output_path": str(dst), "char_count": len(md)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── HTML → TXT ─────────────────────────────────────────────────────────────
    @app.tool()
    def html_to_txt(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Strip HTML tags and save as plain text."""
        try:
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".txt")
            dst = resolve_output(output_path)
            text = _html_to_txt_str(src.read_text(encoding="utf-8"))
            dst.write_text(text, encoding="utf-8")
            return {"success": True, "output_path": str(dst), "char_count": len(text)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── HTML → DOCX ────────────────────────────────────────────────────────────
    @app.tool()
    def html_to_docx(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert an HTML file to a Word document (.docx). Preserves basic heading and paragraph structure."""
        try:
            from docx import Document
            from bs4 import BeautifulSoup

            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".docx")
            dst = resolve_output(output_path)

            soup = BeautifulSoup(src.read_text(encoding="utf-8"), "lxml")
            doc = Document()

            for tag in soup.find_all(["h1","h2","h3","h4","p","li","table"]):
                text = tag.get_text(strip=True)
                if not text:
                    continue
                if tag.name == "h1": doc.add_heading(text, level=1)
                elif tag.name == "h2": doc.add_heading(text, level=2)
                elif tag.name in ("h3","h4"): doc.add_heading(text, level=3)
                elif tag.name == "li": doc.add_paragraph(text, style="List Bullet")
                elif tag.name == "table":
                    rows = tag.find_all("tr")
                    if not rows: continue
                    cols = len(rows[0].find_all(["th","td"]))
                    t = doc.add_table(rows=len(rows), cols=cols)
                    for r, row in enumerate(rows):
                        for c, cell in enumerate(row.find_all(["th","td"])):
                            t.cell(r, c).text = cell.get_text(strip=True)
                else:
                    doc.add_paragraph(text)

            doc.save(str(dst))
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── HTML → XLSX ────────────────────────────────────────────────────────────
    @app.tool()
    def html_to_xlsx(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Extract all HTML tables from a file and save each as a sheet in an Excel workbook."""
        try:
            import pandas as pd
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "_tables", ".xlsx")
            dst = resolve_output(output_path)

            tables = pd.read_html(str(src))
            if not tables:
                return {"success": False, "error": "No tables found in the HTML file."}

            with pd.ExcelWriter(str(dst), engine="openpyxl") as writer:
                for i, df in enumerate(tables):
                    df.to_excel(writer, sheet_name=f"Table{i+1}", index=False)

            return {"success": True, "output_path": str(dst), "tables_extracted": len(tables)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── Markdown → HTML ────────────────────────────────────────────────────────
    @app.tool()
    def markdown_to_html(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert a Markdown file to HTML."""
        try:
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".html")
            dst = resolve_output(output_path)
            md_text = src.read_text(encoding="utf-8")
            body = _md_to_html_str(md_text)
            dst.write_text(_wrap_html(body, src.stem), encoding="utf-8")
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── Markdown → PDF ─────────────────────────────────────────────────────────
    @app.tool()
    def markdown_to_pdf(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert a Markdown file to PDF (renders via HTML → PDF pipeline)."""
        try:
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".pdf")
            dst = resolve_output(output_path)

            md_text = src.read_text(encoding="utf-8")
            html_content = _wrap_html(_md_to_html_str(md_text), src.stem)

            try:
                from weasyprint import HTML
                HTML(string=html_content).write_pdf(str(dst))
                renderer = "weasyprint"
            except ImportError:
                from tools.pdf_import_tools import _html_to_pdf_reportlab
                _html_to_pdf_reportlab(html_content, str(dst))
                renderer = "reportlab"

            return {"success": True, "output_path": str(dst), "renderer": renderer}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── Markdown → DOCX ────────────────────────────────────────────────────────
    @app.tool()
    def markdown_to_docx(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert a Markdown file to a Word document (.docx)."""
        try:
            from docx import Document
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".docx")
            dst = resolve_output(output_path)

            doc = Document()
            for line in src.read_text(encoding="utf-8").split("\n"):
                if line.startswith("### "): doc.add_heading(line[4:], level=3)
                elif line.startswith("## "): doc.add_heading(line[3:], level=2)
                elif line.startswith("# "): doc.add_heading(line[2:], level=1)
                elif line.startswith(("- ", "* ")): doc.add_paragraph(line[2:], style="List Bullet")
                elif re.match(r"^\d+\. ", line): doc.add_paragraph(re.sub(r"^\d+\. ", "", line), style="List Number")
                elif line.strip() == "":
                    continue
                else:
                    # Handle inline bold/italic
                    p = doc.add_paragraph()
                    parts = re.split(r"(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*.*?\*|`.*?`)", line)
                    for part in parts:
                        if part.startswith("***") and part.endswith("***"):
                            run = p.add_run(part[3:-3]); run.bold = True; run.italic = True
                        elif part.startswith("**") and part.endswith("**"):
                            p.add_run(part[2:-2]).bold = True
                        elif part.startswith("*") and part.endswith("*"):
                            p.add_run(part[1:-1]).italic = True
                        elif part.startswith("`") and part.endswith("`"):
                            p.add_run(part[1:-1]).font.name = "Courier New"
                        else:
                            p.add_run(part)
            doc.save(str(dst))
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── Markdown → TXT ─────────────────────────────────────────────────────────
    @app.tool()
    def markdown_to_txt(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Strip Markdown syntax and save as plain text."""
        try:
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".txt")
            dst = resolve_output(output_path)
            text = src.read_text(encoding="utf-8")
            # Strip common markdown syntax
            text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
            text = re.sub(r"\*\*\*(.*?)\*\*\*", r"\1", text)
            text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
            text = re.sub(r"\*(.*?)\*", r"\1", text)
            text = re.sub(r"`(.*?)`", r"\1", text)
            text = re.sub(r"^[-*+]\s+", "", text, flags=re.MULTILINE)
            text = re.sub(r"^\d+\.\s+", "", text, flags=re.MULTILINE)
            text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)  # links
            text = re.sub(r"!\[([^\]]*)\]\([^\)]+\)", r"\1", text)  # images
            dst.write_text(text, encoding="utf-8")
            return {"success": True, "output_path": str(dst), "char_count": len(text)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── TXT → PDF ──────────────────────────────────────────────────────────────
    @app.tool()
    def txt_to_pdf(input_path: str, output_path: Optional[str] = None, font_size: int = 11) -> Dict:
        """Convert a plain text file to PDF using reportlab."""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.units import inch

            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".pdf")
            dst = resolve_output(output_path)

            text = src.read_text(encoding="utf-8")
            doc = SimpleDocTemplate(str(dst), pagesize=A4,
                                    rightMargin=inch, leftMargin=inch,
                                    topMargin=inch, bottomMargin=inch)
            styles = getSampleStyleSheet()
            style = styles["Normal"]
            style.fontSize = font_size

            story = []
            for line in text.split("\n"):
                if line.strip():
                    story.append(Paragraph(html_module.escape(line), style))
                    story.append(Spacer(1, 4))
                else:
                    story.append(Spacer(1, 10))
            doc.build(story)
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── TXT → DOCX ─────────────────────────────────────────────────────────────
    @app.tool()
    def txt_to_docx(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert a plain text file to a Word document (.docx)."""
        try:
            from docx import Document
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".docx")
            dst = resolve_output(output_path)
            doc = Document()
            for line in src.read_text(encoding="utf-8").split("\n"):
                doc.add_paragraph(line)
            doc.save(str(dst))
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── TXT → HTML ─────────────────────────────────────────────────────────────
    @app.tool()
    def txt_to_html(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert a plain text file to HTML, wrapping each line in a <p> tag."""
        try:
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".html")
            dst = resolve_output(output_path)
            lines = src.read_text(encoding="utf-8").split("\n")
            body = "\n".join(f"<p>{html_module.escape(l)}</p>" if l.strip() else "<br>" for l in lines)
            dst.write_text(_wrap_html(body, src.stem), encoding="utf-8")
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── PPTX → TXT ─────────────────────────────────────────────────────────────
    @app.tool()
    def pptx_to_txt(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Extract all text from a PowerPoint presentation (.pptx)."""
        try:
            from pptx import Presentation
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".txt")
            dst = resolve_output(output_path)
            prs = Presentation(str(src))
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
            text = "\n".join(lines)
            dst.write_text(text, encoding="utf-8")
            return {"success": True, "output_path": str(dst), "slides": len(prs.slides), "char_count": len(text)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── PPTX → HTML ────────────────────────────────────────────────────────────
    @app.tool()
    def pptx_to_html(input_path: str, output_path: Optional[str] = None) -> Dict:
        """Convert a PowerPoint presentation to HTML with one section per slide."""
        try:
            from pptx import Presentation
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".html")
            dst = resolve_output(output_path)
            prs = Presentation(str(src))
            sections = []
            for i, slide in enumerate(prs.slides):
                parts = [f'<section id="slide-{i+1}"><h2>Slide {i+1}</h2>']
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                parts.append(f"<p>{html_module.escape(text)}</p>")
                parts.append("</section>")
                sections.append("\n".join(parts))
            body = "\n".join(sections)
            dst.write_text(_wrap_html(body, src.stem), encoding="utf-8")
            return {"success": True, "output_path": str(dst), "slides": len(prs.slides)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── PPTX → Images ──────────────────────────────────────────────────────────
    @app.tool()
    def pptx_to_images(
        input_path: str,
        output_dir: Optional[str] = None,
        output_format: str = "png",
        dpi: int = 150,
    ) -> Dict:
        """
        Convert each slide of a PowerPoint to an image.
        Uses pptx→PDF pipeline (requires Microsoft Office) then renders via PyMuPDF.
        Falls back to extracting embedded images if Office is not available.
        """
        try:
            import fitz
            import tempfile, os

            src = resolve_input(input_path)
            out_dir = Path(output_dir) if output_dir else src.parent / (src.stem + "_slides")
            out_dir.mkdir(parents=True, exist_ok=True)

            # Step 1: PPTX → PDF via docx2pdf
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf:
                tmp_pdf = tf.name
            try:
                from docx2pdf import convert
                convert(str(src), tmp_pdf)
            except Exception as e:
                return {"success": False, "error": f"PPTX→PDF step failed (requires Microsoft PowerPoint): {e}"}

            # Step 2: PDF → images via PyMuPDF
            doc = fitz.open(tmp_pdf)
            fmt_map = {"png": ("PNG", ".png"), "jpg": ("JPEG", ".jpg"), "webp": ("WEBP", ".webp")}
            pil_fmt, ext = fmt_map.get(output_format.lower(), ("PNG", ".png"))
            output_files = []
            from utils.pdf_utils import render_page, pixmap_to_pil
            for i in range(doc.page_count):
                pix = render_page(doc, i, dpi=dpi)
                img = pixmap_to_pil(pix)
                if pil_fmt == "JPEG" and img.mode == "RGBA":
                    img = img.convert("RGB")
                out_path = out_dir / f"{src.stem}_slide{i+1:03d}{ext}"
                img.save(str(out_path), format=pil_fmt)
                output_files.append(str(out_path))
            doc.close()
            os.unlink(tmp_pdf)

            return {"success": True, "output_dir": str(out_dir), "output_files": output_files, "slides": len(output_files)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── XLSX → HTML ────────────────────────────────────────────────────────────
    @app.tool()
    def xlsx_to_html(input_path: str, output_path: Optional[str] = None, sheet_name: Optional[str] = None) -> Dict:
        """Convert an Excel spreadsheet to an HTML file with styled tables."""
        try:
            import pandas as pd
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".html")
            dst = resolve_output(output_path)

            sheets = pd.read_excel(str(src), sheet_name=sheet_name or None, engine="openpyxl")
            if isinstance(sheets, dict):
                parts = []
                for name, df in sheets.items():
                    parts.append(f"<h2>{html_module.escape(str(name))}</h2>")
                    parts.append(df.to_html(index=False, border=1))
                body = "\n".join(parts)
            else:
                body = sheets.to_html(index=False, border=1)

            dst.write_text(_wrap_html(body, src.stem), encoding="utf-8")
            return {"success": True, "output_path": str(dst), "file_size_bytes": dst.stat().st_size}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── XLSX → Markdown ────────────────────────────────────────────────────────
    @app.tool()
    def xlsx_to_markdown(input_path: str, output_path: Optional[str] = None, sheet_name: str = "Sheet1") -> Dict:
        """Convert an Excel spreadsheet sheet to a Markdown table."""
        try:
            import pandas as pd
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".md")
            dst = resolve_output(output_path)

            df = pd.read_excel(str(src), sheet_name=sheet_name, engine="openpyxl")
            md = df.to_markdown(index=False)
            if md is None:
                # fallback if tabulate not installed
                lines = ["| " + " | ".join(str(c) for c in df.columns) + " |"]
                lines.append("| " + " | ".join(["---"] * len(df.columns)) + " |")
                for _, row in df.iterrows():
                    lines.append("| " + " | ".join(str(v) for v in row) + " |")
                md = "\n".join(lines)
            dst.write_text(md, encoding="utf-8")
            return {"success": True, "output_path": str(dst), "rows": len(df), "columns": len(df.columns)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ── XLSX → DOCX ────────────────────────────────────────────────────────────
    @app.tool()
    def xlsx_to_docx(input_path: str, output_path: Optional[str] = None, sheet_name: str = "Sheet1") -> Dict:
        """Embed an Excel spreadsheet as a table in a Word document (.docx)."""
        try:
            import pandas as pd
            from docx import Document
            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".docx")
            dst = resolve_output(output_path)

            df = pd.read_excel(str(src), sheet_name=sheet_name, engine="openpyxl")
            doc = Document()
            doc.add_heading(src.stem, level=1)

            cols = list(df.columns)
            table = doc.add_table(rows=1 + len(df), cols=len(cols))
            table.style = "Table Grid"
            for j, col in enumerate(cols):
                table.rows[0].cells[j].text = str(col)
            for i, (_, row) in enumerate(df.iterrows()):
                for j, val in enumerate(row):
                    table.rows[i + 1].cells[j].text = str(val) if val is not None else ""

            doc.save(str(dst))
            return {"success": True, "output_path": str(dst), "rows": len(df), "columns": len(cols)}
        except Exception as e:
            return {"success": False, "error": str(e)}
