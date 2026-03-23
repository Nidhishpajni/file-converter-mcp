"""PDF export tools: convert PDF to images, Word, Excel, PPTX, HTML, and plain text."""
import io
import os
from pathlib import Path
from typing import Dict, List, Optional

import fitz  # PyMuPDF
from mcp.server.fastmcp import FastMCP

from utils.file_utils import resolve_input, resolve_output, make_output_path
from utils.pdf_utils import open_pdf, render_page, pixmap_to_pil


def register_pdf_export_tools(app: FastMCP) -> None:

    @app.tool()
    def pdf_to_images(
        input_path: str,
        output_dir: Optional[str] = None,
        output_format: str = "png",
        dpi: int = 150,
        pages: Optional[List[int]] = None,
    ) -> Dict:
        """
        Render each PDF page as an image file.

        output_format: png, jpg, or webp (default: png).
        dpi: rendering resolution (default 150; use 300 for print quality).
        pages: 1-based page numbers to render (default: all pages).
        """
        from PIL import Image

        _FMT = {"png": ("PNG", ".png"), "jpg": ("JPEG", ".jpg"), "jpeg": ("JPEG", ".jpg"), "webp": ("WEBP", ".webp")}
        fmt_key = output_format.lower()
        if fmt_key not in _FMT:
            return {"success": False, "error": f"Unsupported format '{output_format}'. Use png, jpg, or webp."}

        pil_fmt, ext = _FMT[fmt_key]

        try:
            src = resolve_input(input_path)
            out_dir = Path(output_dir) if output_dir else src.parent / (src.stem + "_images")
            out_dir.mkdir(parents=True, exist_ok=True)

            doc = open_pdf(str(src))
            n = doc.page_count
            targets = [p - 1 for p in pages] if pages else list(range(n))
            targets = [i for i in targets if 0 <= i < n]

            output_files = []
            for i in targets:
                pix = render_page(doc, i, dpi=dpi)
                img = pixmap_to_pil(pix)

                if pil_fmt == "JPEG" and img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")

                fname = f"{src.stem}_page{i+1:04d}{ext}"
                out_path = out_dir / fname
                kwargs = {"quality": 90, "optimize": True} if pil_fmt in ("JPEG", "WEBP") else {"optimize": True} if pil_fmt == "PNG" else {}
                img.save(str(out_path), format=pil_fmt, **kwargs)
                output_files.append(str(out_path))

            doc.close()
            return {
                "success": True,
                "output_dir": str(out_dir),
                "output_files": output_files,
                "pages_rendered": len(output_files),
                "dpi": dpi,
                "format": pil_fmt,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @app.tool()
    def pdf_to_text(
        input_path: str,
        output_path: Optional[str] = None,
        page_separator: str = "\n\n--- Page {page} ---\n\n",
        password: Optional[str] = None,
    ) -> Dict:
        """
        Extract plain text from a PDF.

        page_separator: template string inserted between pages. Use {page} for page number.
        For scanned PDFs with no text layer, use ocr_pdf instead.
        """
        try:
            import pdfplumber

            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "_text", ".txt")
            dst = resolve_output(output_path)

            open_kwargs = {"password": password} if password else {}
            lines = []
            char_count = 0

            with pdfplumber.open(str(src), **open_kwargs) as pdf:
                for i, page in enumerate(pdf.pages):
                    if i > 0:
                        lines.append(page_separator.format(page=i + 1))
                    text = page.extract_text() or ""
                    lines.append(text)
                    char_count += len(text)

            content = "".join(lines)
            dst.write_text(content, encoding="utf-8")

            return {
                "success": True,
                "output_path": str(dst),
                "char_count": char_count,
                "note": "If output is empty, the PDF may be scanned. Use ocr_pdf for scanned documents.",
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @app.tool()
    def pdf_to_word(
        input_path: str,
        output_path: Optional[str] = None,
        password: Optional[str] = None,
    ) -> Dict:
        """
        Convert a PDF to a Word (.docx) document.

        Extracts text with basic formatting (headings, paragraphs).
        Note: Complex layouts, tables, and images may not transfer perfectly.
        For best results use a text-based PDF (not scanned).
        """
        try:
            import pdfplumber
            from docx import Document
            from docx.shared import Pt

            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "_converted", ".docx")
            dst = resolve_output(output_path)

            doc = Document()
            char_count = 0
            open_kwargs = {"password": password} if password else {}

            with pdfplumber.open(str(src), **open_kwargs) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    if page_num > 0:
                        doc.add_page_break()

                    text = page.extract_text() or ""
                    char_count += len(text)

                    for line in text.split("\n"):
                        line = line.strip()
                        if not line:
                            continue
                        if len(line) < 60 and line.isupper():
                            doc.add_heading(line, level=2)
                        else:
                            doc.add_paragraph(line)

            doc.save(str(dst))

            return {
                "success": True,
                "output_path": str(dst),
                "char_count": char_count,
                "note": "Basic text extraction. Complex layouts may need manual adjustment.",
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @app.tool()
    def pdf_to_excel(
        input_path: str,
        output_path: Optional[str] = None,
        password: Optional[str] = None,
    ) -> Dict:
        """
        Extract tables from a PDF and save them to an Excel (.xlsx) file.

        Each table is placed on a separate sheet named 'Page{n}_Table{t}'.
        Works best on PDFs with clear tabular structure.
        """
        try:
            import pdfplumber
            import openpyxl

            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "_tables", ".xlsx")
            dst = resolve_output(output_path)

            wb = openpyxl.Workbook()
            wb.remove(wb.active)

            tables_found = 0
            open_kwargs = {"password": password} if password else {}

            with pdfplumber.open(str(src), **open_kwargs) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    tables = page.extract_tables() or []
                    for t_idx, table in enumerate(tables):
                        sheet_name = f"Page{page_num+1}_T{t_idx+1}"
                        ws = wb.create_sheet(title=sheet_name)
                        for row in table:
                            ws.append([cell if cell is not None else "" for cell in row])
                        tables_found += 1

            if tables_found == 0:
                wb.create_sheet("Text")
                ws = wb["Text"]
                with pdfplumber.open(str(src), **open_kwargs) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text() or ""
                        for line in text.split("\n"):
                            ws.append([line])

            wb.save(str(dst))

            return {
                "success": True,
                "output_path": str(dst),
                "tables_extracted": tables_found,
                "note": "If 0 tables found, raw text was written as fallback.",
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @app.tool()
    def pdf_to_pptx(
        input_path: str,
        output_path: Optional[str] = None,
        mode: str = "image",
        dpi: int = 150,
        password: Optional[str] = None,
    ) -> Dict:
        """
        Convert a PDF to a PowerPoint (.pptx) presentation.

        mode:
          'image' — each page rendered as a full-slide image (accurate, not searchable).
          'text'  — text extracted and placed as text boxes (searchable, less accurate).
        dpi: rendering resolution for 'image' mode (default 150).
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import pdfplumber

            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "_converted", ".pptx")
            dst = resolve_output(output_path)

            doc = open_pdf(str(src), password=password)
            n = doc.page_count

            prs = Presentation()
            blank_layout = prs.slide_layouts[6]

            if mode == "image":
                for i in range(n):
                    pix = render_page(doc, i, dpi=dpi)
                    img = pixmap_to_pil(pix)
                    if img.mode == "RGBA":
                        img = img.convert("RGB")

                    img_buf = io.BytesIO()
                    img.save(img_buf, format="JPEG", quality=90)
                    img_buf.seek(0)

                    slide = prs.slides.add_slide(blank_layout)
                    w_in = pix.width / dpi
                    h_in = pix.height / dpi
                    prs.slide_width = int(w_in * 914400)
                    prs.slide_height = int(h_in * 914400)

                    slide.shapes.add_picture(img_buf, 0, 0, prs.slide_width, prs.slide_height)

            else:
                open_kwargs = {"password": password} if password else {}
                with pdfplumber.open(str(src), **open_kwargs) as pdf_pl:
                    for page_num, page in enumerate(pdf_pl.pages):
                        slide = prs.slides.add_slide(blank_layout)
                        text = page.extract_text() or ""
                        txBox = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.5),
                            Inches(8.5), Inches(6.5)
                        )
                        txBox.text_frame.text = text

            doc.close()
            prs.save(str(dst))

            return {
                "success": True,
                "output_path": str(dst),
                "slides_created": n,
                "mode": mode,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @app.tool()
    def pdf_to_html(
        input_path: str,
        output_path: Optional[str] = None,
        password: Optional[str] = None,
    ) -> Dict:
        """Convert a PDF to an HTML file with structured text content."""
        import html as html_module

        try:
            import pdfplumber

            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "_converted", ".html")
            dst = resolve_output(output_path)

            open_kwargs = {"password": password} if password else {}
            pages_html = []
            char_count = 0

            with pdfplumber.open(str(src), **open_kwargs) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    char_count += len(text)
                    lines_html = []
                    for line in text.split("\n"):
                        line = line.strip()
                        if not line:
                            continue
                        escaped = html_module.escape(line)
                        if len(line) < 60 and line.isupper():
                            lines_html.append(f"  <h2>{escaped}</h2>")
                        else:
                            lines_html.append(f"  <p>{escaped}</p>")
                    pages_html.append(
                        f'<section id="page-{page_num+1}">\n'
                        f'  <h3 class="page-label">Page {page_num+1}</h3>\n'
                        + "\n".join(lines_html) +
                        "\n</section>"
                    )

            full_html = (
                "<!DOCTYPE html>\n<html lang=\"en\">\n<head>\n"
                "  <meta charset=\"UTF-8\">\n"
                f"  <title>{html_module.escape(src.stem)}</title>\n"
                "  <style>body{font-family:sans-serif;max-width:800px;margin:auto;padding:2em}"
                "section{border-bottom:1px solid #ccc;margin-bottom:2em}"
                ".page-label{color:#999;font-size:.8em}</style>\n"
                "</head>\n<body>\n"
                + "\n".join(pages_html) +
                "\n</body>\n</html>"
            )

            dst.write_text(full_html, encoding="utf-8")

            return {
                "success": True,
                "output_path": str(dst),
                "char_count": char_count,
                "pages": len(pages_html),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    @app.tool()
    def pdf_to_markdown(
        input_path: str,
        output_path: Optional[str] = None,
        password: Optional[str] = None,
    ) -> Dict:
        """
        Convert a PDF to Markdown, preserving basic heading structure and tables.
        For scanned PDFs, run ocr_pdf first.
        """
        try:
            import pdfplumber

            src = resolve_input(input_path)
            if not output_path:
                output_path = make_output_path(input_path, "", ".md")
            dst = resolve_output(output_path)

            open_kwargs = {"password": password} if password else {}
            pages_md = []
            char_count = 0

            with pdfplumber.open(str(src), **open_kwargs) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    if page_num > 0:
                        pages_md.append("\n\n---\n")

                    tables = page.extract_tables() or []
                    table_texts: set = set()
                    for table in tables:
                        if not table or not table[0]:
                            continue
                        header = [str(c or "").strip() for c in table[0]]
                        rows_md = ["| " + " | ".join(header) + " |",
                                   "| " + " | ".join(["---"] * len(header)) + " |"]
                        for row in table[1:]:
                            rows_md.append("| " + " | ".join(str(c or "").strip() for c in row) + " |")
                        pages_md.append("\n".join(rows_md))
                        for row in table:
                            table_texts.update(str(c or "").strip() for c in row if c)

                    text = page.extract_text() or ""
                    char_count += len(text)
                    for line in text.split("\n"):
                        line = line.strip()
                        if not line or line in table_texts:
                            continue
                        if len(line) < 60 and line.isupper():
                            pages_md.append(f"\n## {line}\n")
                        else:
                            pages_md.append(line)

            content = "\n".join(pages_md)
            dst.write_text(content, encoding="utf-8")
            return {"success": True, "output_path": str(dst), "char_count": char_count}
        except Exception as e:
            return {"success": False, "error": f"pdf_to_markdown failed: {e}"}
